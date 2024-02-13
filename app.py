from flask import Flask, render_template, request, send_from_directory
import os
from docx import Document
from pptx import Presentation
from io import BytesIO
import base64
import fitz

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        file = request.files['file']
        if file:
            filename = file.filename
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

            if filename.endswith('.docx'):
                text = process_word(file_path)
                return render_template('word_template.html', text=text)
            
            elif filename.endswith('.pdf'):
                text = process_pdf(file_path)
                return render_template('pdf_template.html', text=text, filename=filename)

            elif filename.endswith('.pptx'):
                slides_data = process_pptx(file_path)
                return render_template('pptx_template.html', slides=slides_data)
            
            else:
                text = "Unsupported file format"
                slides_data = [{"text": text, "images": []}]
                return render_template('pptx_template.html', slides=slides_data)
    
    return render_template('index.html')

@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

def process_word(file_path):
    document = Document(file_path)
    text = ""
    for paragraph in document.paragraphs:
        text += paragraph.text + "\n"
    return text

def process_pdf(file_path):
    text = ""
    with fitz.open(file_path) as pdf:
        for page in pdf:
            text += page.get_text()
    return text

def process_pptx(file_path):
    presentation = Presentation(file_path)
    slides_data = []
    for slide in presentation.slides:
        slide_data = {"text": "", "images": []}
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_data["text"] += shape.text + "\n"
            elif shape.shape_type == 13:  # 13 corresponds to image shape type
                image_stream = BytesIO()
                img = shape.image
                img_bytes = img.blob
                image_stream.write(img_bytes)
                image_data = base64.b64encode(image_stream.getvalue()).decode('utf-8')
                slide_data["images"].append(image_data)
        slides_data.append(slide_data)
    return slides_data

if __name__ == '__main__':
    app.run(debug=True)


# git zalupa 



# giiiiiiiiiihih
# hui
    


# sdjnkghsadjkgads
# sdlgjsdfl;g/d
# adsljngasd